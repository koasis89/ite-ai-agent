#!/usr/bin/env python3
"""제안서 금지(애매모호) 표현 검사기.

RFP 작성지침은 통상 '할 수 있을 것 같다', '고려하고 있다', '가능하다' 등
애매모호한 표현을 금지하며 평가 시 불가능한 것으로 처리한다.
텍스트/마크다운 파일 또는 stdin을 검사하여 위반 후보를 줄 번호와 함께 출력한다.

사용법:
    python check_expressions.py file1.md file2.txt ...
    python check_expressions.py --pptx proposal.pptx   (pptx 텍스트 추출 후 검사)
종료코드: 위반 0건이면 0, 있으면 1
"""
import re
import sys

# (패턴, 설명, 교체 제안)
PATTERNS = [
    (r"할\s*수\s*있(다|음|습니다|을\s*것)", "가능성 표현", "'~한다 / ~를 제공한다'로 단정"),
    (r"제공할\s*수\s*있", "가능성 표현", "'제공한다'"),
    (r"지원할\s*수\s*있", "가능성 표현", "'지원한다'"),
    (r"가능하다|가능합니다|가능함", "가능성 표현", "'~한다' 단정형"),
    (r"고려하고\s*있|고려\s*중|고려하겠", "미확정 표현", "'~를 적용한다 / 반영한다'"),
    (r"검토하겠|검토\s*중|검토할\s*예정", "미확정 표현", "'~한다'로 확정하거나 삭제"),
    (r"노력하겠|노력할\s*것", "의지 표현", "구체적 행위·수치로 대체"),
    (r"~?할\s*것\s*같|것으로\s*보인다|것으로\s*예상", "추측 표현", "근거를 제시하고 단정"),
    (r"최대한|가급적|되도록", "정도 모호", "수치·기준으로 대체"),
    (r"필요\s*시\s*협의|추후\s*협의|별도\s*협의", "책임 회피성", "절차·시점을 명시"),
    (r"예정이다|예정입니다|예정임", "미확정 표현", "'~한다'"),
]

def extract_pptx_text(path):
    from pptx import Presentation
    lines = []
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs)
                    if t.strip():
                        lines.append((f"slide{i}", t))
    return lines

def check_lines(lines, label):
    """lines: list of (line_no_or_tag, text)"""
    hits = 0
    for loc, text in lines:
        for pat, kind, fix in PATTERNS:
            for m in re.finditer(pat, text):
                hits += 1
                snippet = text.strip()
                if len(snippet) > 60:
                    s = max(0, m.start() - 25)
                    snippet = "…" + text[s:s + 60].strip() + "…"
                print(f"[{label}:{loc}] ({kind}) \"{snippet}\"  → {fix}")
    return hits

def main(argv):
    total = 0
    args = [a for a in argv if a != "--pptx"]
    use_pptx = "--pptx" in argv
    if not args:
        lines = [(str(i + 1), l) for i, l in enumerate(sys.stdin.read().splitlines())]
        total += check_lines(lines, "stdin")
    for path in args:
        if use_pptx or path.lower().endswith(".pptx"):
            lines = extract_pptx_text(path)
        else:
            with open(path, encoding="utf-8") as f:
                lines = [(str(i + 1), l) for i, l in enumerate(f.read().splitlines())]
        total += check_lines(lines, path)
    print(f"\n총 {total}건 검출" + ("" if total else " — 통과"))
    return 1 if total else 0

if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
